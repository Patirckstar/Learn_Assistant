"""多格式文档解析工具"""

import json
from pathlib import Path
from typing import Optional


def parse_file(file_path: str) -> Optional[str]:
    """
    根据文件类型解析文档，返回提取的文本内容。
    支持格式: pdf, docx, txt, md, ppt, json
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    suffix = path.suffix.lower()

    parsers = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".txt": _parse_text,
        ".md": _parse_text,
        ".ppt": _parse_ppt,
        ".pptx": _parse_ppt,
        ".json": _parse_json,
    }

    parser = parsers.get(suffix)
    if parser is None:
        raise ValueError(f"不支持的文件格式: {suffix}，支持的格式: {list(parsers.keys())}")

    return parser(file_path)


def _parse_pdf(file_path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    return "\n".join(texts)


def _parse_docx(file_path: str) -> str:
    from docx import Document

    doc = Document(file_path)
    texts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    return "\n".join(texts)


def _parse_text(file_path: str) -> str:
    # 优先 UTF-8，失败后尝试常见中文编码
    encodings = ["utf-8", "gbk", "gb2312", "gb18030", "latin-1"]
    for enc in encodings:
        try:
            with open(file_path, encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    # 最后 fallback: 忽略无法解码的字节
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        return f.read()


def _parse_ppt(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
    return "\n".join(texts)


def _parse_json(file_path: str) -> str:
    with open(file_path, encoding="utf-8") as f:
        data = json.load(f)

    # JSON 可能为不同结构，转为字符串
    if isinstance(data, str):
        return data
    elif isinstance(data, list):
        return "\n".join(str(item) for item in data)
    elif isinstance(data, dict):
        # 若为键值对结构，拼接 key: value
        lines = []
        for key, value in data.items():
            if isinstance(value, str):
                lines.append(f"{key}: {value}")
            else:
                lines.append(f"{key}: {json.dumps(value, ensure_ascii=False)}")
        return "\n".join(lines)
    else:
        return json.dumps(data, ensure_ascii=False)
